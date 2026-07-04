from __future__ import annotations

import csv
import json
import zipfile
from pathlib import Path

import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

from .schemas import ParsedDocument, ParsedFragment
from .utils import detect_language, normalize_text


LEGACY_METADATA_ONLY = {".doc", ".docm"}
ARCHIVE_STAGED = {".zip", ".rar"}


def parse_document(path: Path) -> ParsedDocument:
    extension = path.suffix.lower()
    if extension == ".pdf":
        return _parse_pdf(path)
    if extension == ".docx":
        return _parse_docx(path)
    if extension in {".xlsx", ".xls"}:
        return _parse_spreadsheet(path)
    if extension == ".pptx":
        return _parse_pptx(path)
    if extension in {".json", ".csv", ".txt", ".md"}:
        return _parse_plain(path)
    if extension in LEGACY_METADATA_ONLY:
        return _metadata_only(path, f"Legacy Office format {extension} is registered in metadata-only mode for v1.", mode="metadata_only")
    if extension in ARCHIVE_STAGED:
        return _parse_archive_metadata(path)
    return _metadata_only(path, f"Unsupported rich parsing for extension {extension}. Metadata-only mode applied.", mode="metadata_only")


def _parse_pdf(path: Path) -> ParsedDocument:
    reader = PdfReader(str(path))
    fragments: list[ParsedFragment] = []
    warnings: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = normalize_text(page.extract_text() or "")
        except Exception as exc:  # pragma: no cover - defensive
            warnings.append(f"Failed to parse page {index}: {exc}")
            text = ""
        if text:
            fragments.append(ParsedFragment(fragment_type="page", text=text, page_number=index, ordinal=index))
    language = detect_language(" ".join(fragment.text for fragment in fragments[:3])) if fragments else "ru"
    metadata = {"parser": "pdf"}
    if reader.pages and not fragments:
        warnings.append("PDF contains pages but yielded no text; OCR is likely required.")
        metadata["mode"] = "needs_ocr"
    return ParsedDocument(title=path.stem, language=language, page_count=len(reader.pages), fragments=fragments, warnings=warnings, metadata=metadata)


def _parse_docx(path: Path) -> ParsedDocument:
    document = Document(str(path))
    fragments: list[ParsedFragment] = []
    ordinal = 0
    for paragraph in document.paragraphs:
        text = normalize_text(paragraph.text)
        if not text:
            continue
        ordinal += 1
        fragments.append(ParsedFragment(fragment_type="paragraph", text=text, ordinal=ordinal))
    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            rows.append(" | ".join(normalize_text(cell.text) for cell in row.cells))
        text = normalize_text("\n".join(rows))
        if text:
            ordinal += 1
            fragments.append(ParsedFragment(fragment_type="table", text=text, ordinal=ordinal, metadata={"table_index": table_index}))
    language = detect_language(" ".join(fragment.text for fragment in fragments[:5])) if fragments else "ru"
    return ParsedDocument(title=path.stem, language=language, fragments=fragments, warnings=[], metadata={"parser": "docx"})


def _parse_spreadsheet(path: Path) -> ParsedDocument:
    fragments: list[ParsedFragment] = []
    excel = pd.ExcelFile(path)
    ordinal = 0
    for sheet_name in excel.sheet_names:
        frame = excel.parse(sheet_name).fillna("")
        rows = frame.astype(str).values.tolist()
        for row_index, row in enumerate(rows, start=1):
            text = normalize_text(" | ".join(row))
            if not text:
                continue
            ordinal += 1
            fragments.append(ParsedFragment(fragment_type="table_row", text=text, ordinal=ordinal, metadata={"sheet_name": sheet_name, "row_index": row_index}))
    language = detect_language(" ".join(fragment.text for fragment in fragments[:10])) if fragments else "ru"
    return ParsedDocument(title=path.stem, language=language, fragments=fragments, warnings=[], metadata={"parser": "spreadsheet", "sheet_count": len(excel.sheet_names)})


def _parse_pptx(path: Path) -> ParsedDocument:
    presentation = Presentation(str(path))
    fragments: list[ParsedFragment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                slide_texts.append(normalize_text(text))
        joined = normalize_text(" ".join(slide_texts))
        if joined:
            fragments.append(ParsedFragment(fragment_type="slide", text=joined, page_number=index, ordinal=index))
    language = detect_language(" ".join(fragment.text for fragment in fragments[:5])) if fragments else "ru"
    return ParsedDocument(title=path.stem, language=language, page_count=len(presentation.slides), fragments=fragments, warnings=[], metadata={"parser": "pptx"})


def _parse_plain(path: Path) -> ParsedDocument:
    extension = path.suffix.lower()
    if extension == ".json":
        return _parse_json(path)
    if extension == ".csv":
        rows = []
        with path.open("r", encoding="utf-8") as fh:
            reader = csv.reader(fh)
            for row in reader:
                rows.append(" | ".join(row))
        text = normalize_text("\n".join(rows))
    else:
        text = normalize_text(path.read_text(encoding="utf-8"))
    language = detect_language(text) if text else "ru"
    return ParsedDocument(title=path.stem, language=language, fragments=[ParsedFragment(fragment_type="text", text=text, ordinal=1)] if text else [], warnings=[], metadata={"parser": "plain"})


def _parse_json(path: Path) -> ParsedDocument:
    data = json.loads(path.read_text(encoding="utf-8"))
    fragments: list[ParsedFragment] = []
    ordinal = 0

    def add(fragment_type: str, text: str, metadata: dict | None = None) -> None:
        nonlocal ordinal
        normalized = normalize_text(text)
        if not normalized:
            return
        ordinal += 1
        fragments.append(ParsedFragment(fragment_type=fragment_type, text=normalized, ordinal=ordinal, metadata=metadata or {}))

    if isinstance(data, dict) and "experts" in data:
        for item in data.get("experts", []):
            add(
                "expert_card",
                f"name: {item.get('name', '')}; organization: {item.get('organization', '')}; topics: {', '.join(item.get('topics', []))}",
                {"entry_type": "expert"},
            )
    elif isinstance(data, dict) and any(key in data for key in {"materials", "equipment", "units"}):
        for group_name in ["materials", "equipment", "units"]:
            for item in data.get(group_name, []):
                aliases = item.get("aliases", [])
                parts = [f"group: {group_name}", f"name: {item.get('name', '')}"]
                if aliases:
                    parts.append(f"aliases: {', '.join(aliases)}")
                if item.get("property"):
                    parts.append(f"property: {item.get('property')}")
                if item.get("normalized"):
                    parts.append(f"normalized: {item.get('normalized')}")
                add("reference_entry", "; ".join(parts), {"entry_type": group_name})
    elif isinstance(data, dict) and "topics" in data:
        for item in data.get("topics", []):
            add(
                "taxonomy_topic",
                f"topic: {item.get('name', '')}; parent: {item.get('parent') or ''}",
                {"entry_type": "topic"},
            )
    elif isinstance(data, dict) and any(key in data for key in {"material", "process", "regime", "result"}):
        if data.get("title"):
            add("experiment_title", f"title: {data.get('title')}", {"entry_type": "experiment_title"})
        if data.get("material"):
            add("experiment_material", f"material: {data.get('material')}", {"entry_type": "material"})
        if data.get("process"):
            add("experiment_process", f"process: {data.get('process')}", {"entry_type": "process"})
        for key, value in (data.get("regime") or {}).items():
            add("experiment_regime", f"{key}: {value}", {"entry_type": "regime", "key": key})
        for key, value in (data.get("result") or {}).items():
            add("experiment_result", f"{key}: {value}", {"entry_type": "result", "key": key})
        if data.get("notes"):
            add("experiment_notes", f"notes: {data.get('notes')}", {"entry_type": "notes"})
    else:
        add("json", json.dumps(data, ensure_ascii=False), {"entry_type": "generic_json"})

    full_text = normalize_text(json.dumps(data, ensure_ascii=False))
    if full_text:
        add("json_blob", full_text, {"entry_type": "json_blob"})

    language = detect_language(" ".join(fragment.text for fragment in fragments[:8])) if fragments else "ru"
    return ParsedDocument(title=path.stem, language=language, fragments=fragments, warnings=[], metadata={"parser": "json", "structured": True})


def _parse_archive_metadata(path: Path) -> ParsedDocument:
    members: list[str] = []
    warnings: list[str] = []
    if path.suffix.lower() == ".zip":
        try:
            with zipfile.ZipFile(path) as archive:
                members = sorted(item.filename for item in archive.infolist() if not item.is_dir())[:100]
        except Exception as exc:  # pragma: no cover - defensive
            warnings.append(f"Failed to inspect archive members: {exc}")
    else:
        warnings.append("RAR staged import is registered but member inspection is not enabled in v1 runtime.")
    warnings.append("Archive source was registered as staged import. Extract supported files and re-import for deep parsing.")
    metadata = {"mode": "staged_import", "archive_members": members, "parser": "archive"}
    return ParsedDocument(title=path.stem, language="ru", fragments=[], warnings=warnings, metadata=metadata)


def _metadata_only(path: Path, warning: str, *, mode: str) -> ParsedDocument:
    return ParsedDocument(title=path.stem, language="ru", fragments=[], warnings=[warning], metadata={"mode": mode})
